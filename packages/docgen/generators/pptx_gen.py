"""PowerPoint generator using python-pptx."""
from __future__ import annotations

import io

from pptx import Presentation
from pptx.util import Inches, Pt

from packages.docgen.generators.base import DocumentContent, DocumentGenerator, DocumentSection


class PptxGenerator(DocumentGenerator):
    """Generate .pptx files from DocumentContent."""

    @property
    def mime_type(self) -> str:
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    @property
    def file_extension(self) -> str:
        return "pptx"

    def generate(
        self, content: DocumentContent, template_path: str | None = None
    ) -> bytes:
        prs = Presentation(template_path) if template_path else Presentation()

        # Title slide
        title_layout = prs.slide_layouts[0]  # Title Slide layout
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = content.title
        if slide.placeholders[1]:
            author = content.metadata.get("author", "")
            slide.placeholders[1].text = author

        # Content slides
        for section in content.sections:
            self._add_section_slide(prs, section)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # ------------------------------------------------------------------
    def _add_section_slide(self, prs: Presentation, section: DocumentSection) -> None:
        # Use Title and Content layout
        content_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_layout)

        slide.shapes.title.text = section.heading

        # Body content
        if section.body and slide.placeholders[1]:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for i, line in enumerate(section.body.split("\n")):
                stripped = line.strip()
                if not stripped:
                    continue
                if i == 0:
                    tf.text = stripped
                else:
                    p = tf.add_paragraph()
                    p.text = stripped

        # Table on slide
        if section.table:
            self._add_table(slide, section.table)

        # Slide notes
        if section.slide_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = section.slide_notes

    def _add_table(self, slide, table_data: list[list[str]]) -> None:
        if not table_data:
            return

        rows = len(table_data)
        cols = max(len(r) for r in table_data)

        left = Inches(0.5)
        top = Inches(4.0)
        width = Inches(9.0)
        height = Inches(0.8)

        tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        tbl = tbl_shape.table

        for r_idx, row in enumerate(table_data):
            for c_idx, cell_text in enumerate(row):
                if c_idx < cols:
                    tbl.cell(r_idx, c_idx).text = cell_text
